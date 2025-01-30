import csv
import json
from pathlib import Path
import requests
from bs4 import BeautifulSoup


class CommandHandler:
    """指令处理类"""
    def __init__(self):
        self.running = True

    def handle_command(self, user_input: str) -> [str,str]:
        if user_input.startswith('/exit'):
            self.running = False
            return "正在退出程序...", ""
        elif user_input.startswith('/file'):
            return self.process_file(user_input)
        elif user_input.startswith('/fetch'):  # 新增网页获取指令
            return self.process_fetch(user_input)
        else:
            return f"未知指令: {user_input.split()[0]}", ""

    @staticmethod
    def process_fetch(command: str) -> [str, str]:
        """处理网页获取指令"""
        try:
            _, url = command.split(maxsplit=1)
            url = url.strip()
            return f"已获取网页内容：{url}", f"[网页内容摘要]: {fetch_web_content(url)}"
        except ValueError:
            return "URL参数缺失", ""
        except Exception as e:
            return f"网页获取失败: {str(e)}", ""

    @staticmethod
    def process_file(command: str) -> [str, str]:
        """处理文件读取指令"""
        try:
            _, filepath = command.split(maxsplit=1)
            return f"读取到了本地文件{filepath}", f"[解析的本地文件{filepath}]: {read_file_content(filepath.strip())}"
        except ValueError:
            return "文件路径参数缺失", ""
        except Exception as e:
            return f"文件处理失败: {str(e)}", ""


def fetch_web_content(url: str, timeout: int = 10) -> str:
    """获取并解析网页内容"""
    try:
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }

        # 带流式传输的请求以防止大文件
        with requests.get(url, headers=headers, timeout=timeout, stream=True) as response:
            response.raise_for_status()

            # 检测编码
            if response.encoding is None:
                response.encoding = 'utf-8'

            # 读取部分内容进行解析
            content = []
            for chunk in response.iter_content(decode_unicode=True, chunk_size=1024):
                if chunk:
                    content.append(chunk)
                if len(content) > 20:  # 限制读取20个chunk（约20KB）
                    break

            text = ''.join(content)

            # 使用BeautifulSoup提取主要内容
            soup = BeautifulSoup(text, 'html.parser')

            # 移除不需要的元素
            for element in soup(['script', 'style', 'nav', 'footer', 'header', 'meta']):
                element.decompose()

            # 提取文本内容
            content = []
            for tag in soup.find_all(['p', 'h1', 'h2', 'h3', 'article']):
                content.append(tag.get_text(strip=True, separator=' '))

            # 合并并限制输出长度
            full_text = '\n'.join(content)
            return full_text[:5000] + "..." if len(full_text) > 5000 else full_text

    except requests.exceptions.RequestException as e:
        raise RuntimeError(f"网络请求失败: {str(e)}")
    except ImportError:
        raise ImportError("需要安装依赖库：requests和beautifulsoup4，请执行 `pip install requests beautifulsoup4`")
    except Exception as e:
        raise RuntimeError(f"内容解析失败: {str(e)}")


def read_file_content(filepath: str) -> str:
    """读取并解析多种文件格式"""
    path = Path(filepath)
    if not path.exists():
        raise FileNotFoundError(f"文件不存在: {filepath}")

    suffix = path.suffix.lower()

    if suffix == '.txt':
        return path.read_text(encoding='utf-8')
    elif suffix == '.csv':
        return csv_to_text(path)
    elif suffix == '.json':
        return json_to_text(path)
    elif suffix == '.py':
        return f"Python代码文件内容:\n```python\n{path.read_text(encoding='utf-8')}\n```"
    elif suffix == '.pdf':
        return pdf_to_text(path)
    elif suffix in ('.docx', '.doc'):
        return docx_to_text(path)
    elif suffix in ('.pptx', '.ppt'):
        return pptx_to_text(path)
    elif suffix in ('.xlsx', '.xls'):
        return excel_to_text(path)
    else:
        raise ValueError(f"不支持的文件格式: {suffix}")


def csv_to_text(filepath: Path) -> str:
    """将CSV转换为自然语言描述"""
    with open(filepath, 'r', encoding='utf-8') as f:
        reader = csv.DictReader(f)
        output = [f"CSV文件包含 {len(reader.fieldnames)} 列: {', '.join(reader.fieldnames)}"]
        for i, row in enumerate(reader, 1):
            output.append(f"行 {i}: {', '.join(f'{k}={v}' for k,v in row.items())}")
    return "\n".join(output)


def json_to_text(filepath: Path) -> str:
    """将JSON转换为自然语言描述"""
    data = json.loads(filepath.read_text(encoding='utf-8'))
    return f"JSON数据结构摘要:\n{json.dumps(data, indent=2, ensure_ascii=False)}"


def pdf_to_text(filepath: Path) -> str:
    """提取PDF文件文本内容"""
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        raise ImportError("处理PDF文件需要PyPDF2库，请执行`pip install PyPDF2`安装")

    text = []
    try:
        with open(filepath, 'rb') as f:
            reader = PdfReader(f)
            for page in reader.pages[:10]:  # 限制最多读取10页
                page_text = page.extract_text()
                if page_text:
                    text.append(page_text.strip())
    except Exception as e:
        raise RuntimeError(f"读取PDF文件失败: {str(e)}")

    full_text = "\n".join(text)
    return full_text[:10000] + "..." if len(full_text) > 10000 else full_text


def docx_to_text(filepath: Path) -> str:
    """提取Word文档文本内容"""
    try:
        from docx import Document
    except ImportError:
        raise ImportError("处理Word文件需要python-docx库，请执行`pip install python-docx`安装")

    try:
        doc = Document(filepath)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        raise RuntimeError(f"读取Word文件失败: {str(e)}")


def pptx_to_text(filepath: Path) -> str:
    """提取PPT幻灯片文本内容"""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("处理PPT文件需要python-pptx库，请执行`pip install python-pptx`安装")

    try:
        prs = Presentation(filepath)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text.strip())
        return "\n".join(text)
    except Exception as e:
        raise RuntimeError(f"读取PPT文件失败: {str(e)}")


def excel_to_text(filepath: Path) -> str:
    """提取Excel文件内容摘要"""
    try:
        import pandas as pd
    except ImportError:
        raise ImportError("处理Excel文件需要pandas库，请执行`pip install pandas openpyxl`安装")

    try:
        xls = pd.ExcelFile(filepath)
        output = []
        for sheet_name in xls.sheet_names[:3]:  # 最多处理前3个工作表
            df = xls.parse(sheet_name, nrows=5)  # 每个工作表读取前5行
            output.append(
                f"工作表 '{sheet_name}' 的前5行数据:\n"
                f"{df.to_string(index=False)}\n"
                f"共 {len(df)} 行 × {len(df.columns)} 列"
            )
        return "\n\n".join(output)
    except Exception as e:
        raise RuntimeError(f"读取Excel文件失败: {str(e)}")