import csv
import json
from pathlib import Path
from typing import List, Tuple
from core.cache import CatchInformation

from .commands import registry, Command, CommandContext


@registry.register(
    path="/file",
    description="文件操作",
    usage="/file <文件路径>"
)
class FileCommand(Command):
    def execute(self, args: List[str], context: CommandContext) -> Tuple[str, str]:
        if not args:
            return "文件路径参数缺失", ""

        filepath = args[0].strip()
        try:
            content = read_file_content(filepath)
            CatchInformation.get_instance().info = f"读取到了本地文件 {filepath} 的完整内容：{content}"
            return f"读取到了本地文件 {filepath}， 使用/summary总结关键信息或者使用/submit将内容提交给AI", ""
        except Exception as e:
            return f"文件处理失败: {str(e)}", ""


def read_file_content(filepath: str) -> str:
    """读取并解析多种文件格式"""
    path = Path(filepath)
    if not path.exists():
        raise FileNotFoundError(f"文件不存在: {filepath}")

    suffix = path.suffix.lower()

    if suffix == '.txt':
        return path.read_text(encoding='utf-8')
    elif suffix == '.csv':
        return csv_to_text(path)
    elif suffix == '.json':
        return json_to_text(path)
    elif suffix == '.py':
        return f"带行号的Python代码文件内容:\n```python\n{read_python_code(path)}\n```"
    elif suffix == '.pdf':
        return pdf_to_text(path)
    elif suffix in ('.docx', '.doc'):
        return docx_to_text(path)
    elif suffix in ('.pptx', '.ppt'):
        return pptx_to_text(path)
    elif suffix in ('.xlsx', '.xls'):
        return excel_to_text(path)
    else:
        raise ValueError(f"不支持的文件格式: {suffix}")

def read_python_code(filepath: Path) -> str:
    """读取Python代码文件内容，为每行添加行号"""
    with filepath.open("r", encoding="utf-8") as f:
        lines = f.readlines()

    return "".join(f"{i+1:4d}: {line}" for i, line in enumerate(lines))


def csv_to_text(filepath: Path) -> str:
    """将CSV转换为自然语言描述"""
    with open(filepath, 'r', encoding='utf-8') as f:
        reader = csv.DictReader(f)
        output = [f"CSV文件包含 {len(reader.fieldnames)} 列: {', '.join(reader.fieldnames)}"]
        for i, row in enumerate(reader, 1):
            output.append(f"行 {i}: {', '.join(f'{k}={v}' for k,v in row.items())}")
    return "\n".join(output)


def json_to_text(filepath: Path) -> str:
    """将JSON转换为自然语言描述"""
    data = json.loads(filepath.read_text(encoding='utf-8'))
    return f"JSON数据结构摘要:\n{json.dumps(data, indent=2, ensure_ascii=False)}"


def pdf_to_text(filepath: Path) -> str:
    """提取PDF文件文本内容"""
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        raise ImportError("处理PDF文件需要PyPDF2库，请执行`pip install PyPDF2`安装")

    text = []
    try:
        with open(filepath, 'rb') as f:
            reader = PdfReader(f)
            for page in reader.pages[:10]:  # 限制最多读取10页
                page_text = page.extract_text()
                if page_text:
                    text.append(page_text.strip())
    except Exception as e:
        raise RuntimeError(f"读取PDF文件失败: {str(e)}")

    full_text = "\n".join(text)
    return full_text[:10000] + "..." if len(full_text) > 10000 else full_text


def docx_to_text(filepath: Path) -> str:
    """提取Word文档文本内容"""
    try:
        from docx import Document
    except ImportError:
        raise ImportError("处理Word文件需要python-docx库，请执行`pip install python-docx`安装")

    try:
        doc = Document(filepath)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        raise RuntimeError(f"读取Word文件失败: {str(e)}")


def pptx_to_text(filepath: Path) -> str:
    """提取PPT幻灯片文本内容"""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("处理PPT文件需要python-pptx库，请执行`pip install python-pptx`安装")

    try:
        prs = Presentation(filepath)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text.strip())
        return "\n".join(text)
    except Exception as e:
        raise RuntimeError(f"读取PPT文件失败: {str(e)}")


def excel_to_text(filepath: Path) -> str:
    """提取Excel文件内容摘要"""
    try:
        import pandas as pd
    except ImportError:
        raise ImportError("处理Excel文件需要pandas库，请执行`pip install pandas openpyxl`安装")

    try:
        xls = pd.ExcelFile(filepath)
        output = []
        for sheet_name in xls.sheet_names[:3]:  # 最多处理前3个工作表
            df = xls.parse(sheet_name, nrows=5)  # 每个工作表读取前5行
            output.append(
                f"工作表 '{sheet_name}' 的前5行数据:\n"
                f"{df.to_string(index=False)}\n"
                f"共 {len(df)} 行 × {len(df.columns)} 列"
            )
        return "\n\n".join(output)
    except Exception as e:
        raise RuntimeError(f"读取Excel文件失败: {str(e)}")